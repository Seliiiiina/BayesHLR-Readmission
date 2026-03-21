#!/usr/bin/env python3
"""
Generate presentation PPT for Bayesian Hierarchical Logistic Regression project.
"""

import os
import io
import csv
import textwrap
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

# ── paths ─────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent
FIG_DIR = ROOT / "output" / "figures"
TBL_DIR = ROOT / "output" / "tables"
OUT_PPT = ROOT / "presentation.pptx"
TEMP_DIR = ROOT / "output" / "temp_ppt"
TEMP_DIR.mkdir(parents=True, exist_ok=True)

# ── colours ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x23, 0x3D, 0x6B)
ACCENT    = RGBColor(0x3A, 0x7C, 0xA5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)
TEXT_DARK  = RGBColor(0x2C, 0x2C, 0x2C)
TEXT_GRAY  = RGBColor(0x66, 0x66, 0x66)
HIGHLIGHT  = RGBColor(0xC0, 0x39, 0x2B)
GREEN      = RGBColor(0x1E, 0x82, 0x49)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── helpers ───────────────────────────────────────────────────────────────────

def render_latex(tex, filename, fontsize=20, dpi=200):
    """Render a LaTeX string to a PNG image and return its path."""
    path = TEMP_DIR / filename
    fig, ax = plt.subplots(figsize=(12, 1.2))
    ax.text(0.5, 0.5, tex, fontsize=fontsize, ha="center", va="center",
            transform=ax.transAxes, usetex=False, math_fontfamily="cm")
    ax.axis("off")
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    fig.savefig(str(path), dpi=dpi, bbox_inches="tight", transparent=True, pad_inches=0.1)
    plt.close(fig)
    return str(path)


def render_latex_block(tex_lines, filename, fontsize=18, dpi=200, figheight=None):
    """Render multiple lines of LaTeX."""
    n = len(tex_lines)
    h = figheight or max(1.0, 0.55 * n)
    fig, ax = plt.subplots(figsize=(12, h))
    for i, line in enumerate(tex_lines):
        y = 1 - (i + 0.5) / n
        ax.text(0.5, y, line, fontsize=fontsize, ha="center", va="center",
                transform=ax.transAxes, usetex=False, math_fontfamily="cm")
    ax.axis("off")
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    fig.savefig(str(path := TEMP_DIR / filename), dpi=dpi, bbox_inches="tight",
                transparent=True, pad_inches=0.1)
    plt.close(fig)
    return str(path)


def add_background(slide, color):
    """Set slide background colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    """Add a coloured rectangle behind everything."""
    w = width or SLIDE_W
    h = height or SLIDE_H
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.3):
    """Add a textbox with a single run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_slide(slide, title, bullets, sub_bullets=None, title_size=28,
                     bullet_size=18, start_top=None):
    """Add title bar + bullet points to a slide."""
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
                title, font_size=title_size, bold=True, color=WHITE)

    top = start_top or Inches(1.4)
    for i, bullet in enumerate(bullets):
        # Check if this bullet has sub-bullets
        subs = None
        if sub_bullets and i in sub_bullets:
            subs = sub_bullets[i]

        tb = add_textbox(slide, Inches(0.8), top, Inches(11.5), Inches(0.5),
                         "", font_size=bullet_size, color=TEXT_DARK)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = TEXT_DARK
        run.font.name = "Calibri"
        p.line_spacing = Pt(bullet_size * 1.5)

        top += Inches(0.48)

        if subs:
            for sub in subs:
                stb = add_textbox(slide, Inches(1.4), top, Inches(10.5), Inches(0.45),
                                  "", font_size=bullet_size - 2, color=TEXT_GRAY)
                stf = stb.text_frame
                stf.word_wrap = True
                sp = stf.paragraphs[0]
                sr = sp.add_run()
                sr.text = "  " + sub
                sr.font.size = Pt(bullet_size - 2)
                sr.font.color.rgb = TEXT_GRAY
                sr.font.name = "Calibri"
                top += Inches(0.38)

    return top


def add_title_bar(slide, title, title_size=28):
    """Just add the navy title bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
                title, font_size=title_size, bold=True, color=WHITE)


def section_divider(prs, title, subtitle=""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_shape_bg(slide, DARK_BLUE)

    # decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.5), Inches(3.1), Inches(2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_textbox(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(1.2),
                title, font_size=36, bold=True, color=WHITE, font_name="Calibri")
    if subtitle:
        add_textbox(slide, Inches(1.5), Inches(4.4), Inches(10), Inches(0.8),
                    subtitle, font_size=20, color=ACCENT, font_name="Calibri")
    return slide


def add_image_centered(slide, img_path, top, max_width=Inches(11), max_height=Inches(5)):
    """Add an image centered horizontally, fitting within bounds."""
    from PIL import Image as PILImage
    try:
        with PILImage.open(img_path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = 1200, 600

    aspect = iw / ih
    # fit within max_width x max_height
    w = max_width
    h = int(w / aspect)
    if h > max_height:
        h = max_height
        w = int(h * aspect)

    left = (SLIDE_W - w) // 2
    slide.shapes.add_picture(img_path, left, top, w, h)


# ── build presentation ───────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# =============================================================================
# SLIDE 1: Title
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(sl, NAVY)

# decorative accent bar
sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                     Inches(1.2), Inches(2.0), Inches(3), Inches(0.07)).fill.solid()
sl.shapes[-1].fill.fore_color.rgb = ACCENT
sl.shapes[-1].line.fill.background()

add_textbox(sl, Inches(1.2), Inches(2.3), Inches(11), Inches(1.2),
            "Bayesian Hierarchical Logistic Regression\nfor 30-Day Hospital Readmission",
            font_size=36, bold=True, color=WHITE)
add_textbox(sl, Inches(1.2), Inches(3.8), Inches(11), Inches(0.7),
            "The Role of Medical Specialty Among Diabetic Patients",
            font_size=22, color=ACCENT)
add_textbox(sl, Inches(1.2), Inches(5.0), Inches(11), Inches(0.5),
            "Selina You,  Zhouhan Qian",
            font_size=20, color=WHITE)
add_textbox(sl, Inches(1.2), Inches(5.6), Inches(11), Inches(0.5),
            "Introduction to Applied Bayesian Analysis  |  March 2026",
            font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC))

# =============================================================================
# SLIDE 2: Section — Motivation
# =============================================================================
section_divider(prs, "Part 1", "Motivation & Research Question")

# =============================================================================
# SLIDE 3: Why Study Readmission?
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(sl, "Why Study Readmission?", [
    "$26 B+  annual cost of unplanned readmissions in the US",
    "CMS Hospital Readmissions Reduction Program penalizes hospitals with excess readmissions",
    "30-day readmission rate is a core quality metric tied to Medicare reimbursement",
    "Diabetic patients are a high-risk group: multiple comorbidities, frequent hospitalizations,\n"
    "    cross-specialty care",
    "Key question: does the admitting specialty itself affect readmission risk?"
])

# =============================================================================
# SLIDE 4: Research Question
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Research Question")

# main question box
qbox = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.6))
qbox.fill.solid()
qbox.fill.fore_color.rgb = LIGHT_BLUE_BG
qbox.line.color.rgb = ACCENT
qbox.line.width = Pt(2)

add_textbox(sl, Inches(1.4), Inches(1.8), Inches(10.5), Inches(1.3),
            "How much specialty-level heterogeneity exists in 30-day\n"
            "readmission risk, and can we reliably estimate it despite\n"
            "severe group imbalance?",
            font_size=22, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(0.8), Inches(3.7), Inches(11.5), Inches(3),
            "Key observations driving this question:",
            font_size=20, bold=True, color=TEXT_DARK)

bullets_rq = [
    "Raw readmission rates vary widely across specialties: 4.3% to 32.3%",
    "But sample sizes are severely imbalanced: 288 (Phys Med) to 10,975 (Internal Med)",
    "Direct comparison of raw rates is unreliable under such imbalance",
    "We need a framework that borrows strength across groups while quantifying uncertainty"
]
top = Inches(4.3)
for b in bullets_rq:
    add_textbox(sl, Inches(1.2), top, Inches(11), Inches(0.5),
                "   " + b, font_size=18, color=TEXT_DARK)
    top += Inches(0.5)

# =============================================================================
# SLIDE 5: Why Not Standard Logistic Regression?
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Why Not Standard Logistic Regression?")

approaches = [
    ("Pooled Model (No Specialty)",
     "Ignores group structure entirely\nTreats all specialties as identical",
     HIGHLIGHT),
    ("Fixed-Effect Dummies",
     "Unstable MLE for small groups (n = 288)\nNo shrinkage toward population mean",
     HIGHLIGHT),
    ("Frequentist Mixed Model",
     "Point estimate for variance component\nNo full uncertainty quantification for heterogeneity",
     HIGHLIGHT),
]

top = Inches(1.5)
for title, desc, col in approaches:
    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), top, Inches(11.5), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    box.line.color.rgb = HIGHLIGHT
    box.line.width = Pt(1)

    add_textbox(sl, Inches(1.2), top + Inches(0.08), Inches(3.5), Inches(0.4),
                title, font_size=18, bold=True, color=HIGHLIGHT)
    add_textbox(sl, Inches(1.2), top + Inches(0.5), Inches(10.5), Inches(0.7),
                desc, font_size=16, color=TEXT_DARK)
    top += Inches(1.5)

# solution box
sol_box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), top + Inches(0.2), Inches(11.5), Inches(1.1))
sol_box.fill.solid()
sol_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
sol_box.line.color.rgb = GREEN
sol_box.line.width = Pt(2)

add_textbox(sl, Inches(1.2), top + Inches(0.35), Inches(10.5), Inches(0.7),
            "Bayesian Hierarchical Model:  partial pooling  +  full uncertainty quantification\n"
            "  +  principled handling of group imbalance",
            font_size=19, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 6: Data Overview
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Data Overview")

data_bullets = [
    "UCI Diabetes 130-US Hospitals dataset (1999-2008, 130 hospitals)",
    "After cleaning:  N = 37,964 encounters  |  J = 19 specialties  |  K = 43 covariates",
    "Overall 30-day readmission rate:  ~9.6%  (3,652 readmitted)",
]
top = Inches(1.3)
for b in data_bullets:
    add_textbox(sl, Inches(0.8), top, Inches(12), Inches(0.4),
                "   " + b, font_size=17, color=TEXT_DARK)
    top += Inches(0.42)

# Add figures side by side
fig1 = str(FIG_DIR / "fig1_specialty_distribution.png")
fig3 = str(FIG_DIR / "fig3_raw_readmission_by_specialty.png")
if os.path.exists(fig1):
    sl.shapes.add_picture(fig1, Inches(0.3), Inches(2.8), Inches(6.3), Inches(4.3))
if os.path.exists(fig3):
    sl.shapes.add_picture(fig3, Inches(6.8), Inches(2.8), Inches(6.3), Inches(4.3))

# =============================================================================
# SLIDE 7: Data Processing Pipeline
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Data Processing Pipeline")

steps = [
    ("Raw Data", "101,766 encounters, 50 variables", RGBColor(0xBB, 0xDE, 0xFB)),
    ("Remove Duplicates", "-30,248 duplicate encounter IDs", RGBColor(0xC8, 0xE6, 0xC9)),
    ("Exclude Expired / Hospice", "-1,545 (discharge codes 11,13,14,19,20,21)", RGBColor(0xC8, 0xE6, 0xC9)),
    ("Exclude Unknown Specialty", "-32,008 records (45.7%)", RGBColor(0xC8, 0xE6, 0xC9)),
    ("Drop High-Missing Vars", "weight (96.9%), payer_code (39.6%)", RGBColor(0xC8, 0xE6, 0xC9)),
    ("Specialty Regrouping", "73 raw specialties -> 19 groups (threshold >= 250)", RGBColor(0xFE, 0xF9, 0xE7)),
    ("Feature Engineering", "10 continuous (z-score) + 33 binary/categorical = K=43", RGBColor(0xFE, 0xF9, 0xE7)),
    ("Final Cohort", "N=37,964  |  J=19  |  K=43  |  Outcome: readmit_30", RGBColor(0xBB, 0xDE, 0xFB)),
]

top = Inches(1.4)
box_h = Inches(0.62)
box_w = Inches(11.5)
arrow_h = Inches(0.18)

for i, (label, desc, bg_col) in enumerate(steps):
    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.9), top, box_w, box_h)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_col
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1)

    add_textbox(sl, Inches(1.2), top + Inches(0.05), Inches(3.2), Inches(0.5),
                label, font_size=16, bold=True, color=DARK_BLUE)
    add_textbox(sl, Inches(4.5), top + Inches(0.05), Inches(7.5), Inches(0.5),
                desc, font_size=15, color=TEXT_DARK)

    top += box_h
    if i < len(steps) - 1:
        # arrow
        arrow = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(6.4), top, Inches(0.4), arrow_h)
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()
        top += arrow_h


# =============================================================================
# SLIDE 8: Section — Overall Framework
# =============================================================================
section_divider(prs, "Part 2", "Overall Analysis Framework")

# =============================================================================
# SLIDE 9: Project Pipeline
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Analysis Pipeline Overview")

pipeline_steps = [
    ("1. Data Cleaning & Cohort", "01_clean_data.py", "N = 37,964"),
    ("2. Feature Engineering", "02_feature_engineering.py", "K=43, J=19"),
    ("3. Exploratory Analysis", "03_eda_tables_figures.py", "Tables & Figures"),
    ("4. Model Specification", "Hierarchical logistic regression", "Priors + Likelihood"),
    ("5. Posterior Derivation", "Log-posterior (non-conjugate)", "Reparameterization"),
    ("6. MCMC Sampling", "05_mwg_sampler.py", "4 chains x 30K iter"),
    ("7. Convergence Diagnostics", "Split-Rhat, ESS, trace plots", "All Rhat < 1.05"),
    ("8. Posterior Inference", "Fixed effects, random effects, tau", "Core findings"),
    ("9. Model Evaluation", "DIC, PPC", "Model adequacy"),
    ("10. Sensitivity Analysis", "Alternative priors", "Robustness"),
]

left_col_x = Inches(0.6)
mid_col_x = Inches(4.2)
right_col_x = Inches(8.5)
top = Inches(1.35)
row_h = Inches(0.55)

# headers
for hx, htxt in [(left_col_x, "Step"), (mid_col_x, "Implementation"), (right_col_x, "Key Output")]:
    add_textbox(sl, hx, top, Inches(3.5), Inches(0.4),
                htxt, font_size=16, bold=True, color=NAVY)
top += Inches(0.45)

# divider line
line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(0.6), top, Inches(12), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()
top += Inches(0.15)

for step, impl, output in pipeline_steps:
    add_textbox(sl, left_col_x, top, Inches(3.5), Inches(0.4),
                step, font_size=14, bold=True, color=TEXT_DARK)
    add_textbox(sl, mid_col_x, top, Inches(4), Inches(0.4),
                impl, font_size=14, color=TEXT_GRAY)
    add_textbox(sl, right_col_x, top, Inches(4), Inches(0.4),
                output, font_size=14, color=ACCENT)
    top += row_h

# =============================================================================
# SLIDE 10: Section — Model Structure
# =============================================================================
section_divider(prs, "Part 3", "Model Structure")

# =============================================================================
# SLIDE 11: Likelihood
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Model Specification — Likelihood")

# Render equations
eq1 = render_latex(r"$y_i \mid \alpha,\,\boldsymbol{\beta},\,\mathbf{u} \;\sim\; \mathrm{Bernoulli}(\pi_i)$",
                   "eq_lik1.png", fontsize=24)
eq2 = render_latex(r"$\eta_i \;=\; \alpha \;+\; \mathbf{x}_i^\top \boldsymbol{\beta} \;+\; u_{j[i]}$",
                   "eq_lik2.png", fontsize=24)
eq3 = render_latex(r"$\pi_i \;=\; \frac{\exp(\eta_i)}{1 + \exp(\eta_i)}$",
                   "eq_lik3.png", fontsize=24)

sl.shapes.add_picture(eq1, Inches(1.5), Inches(1.5), Inches(8), Inches(0.7))
sl.shapes.add_picture(eq2, Inches(1.5), Inches(2.5), Inches(8), Inches(0.7))
sl.shapes.add_picture(eq3, Inches(1.5), Inches(3.5), Inches(8), Inches(0.7))

# Conditional independence
ci_box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.7))
ci_box.fill.solid()
ci_box.fill.fore_color.rgb = LIGHT_BLUE_BG
ci_box.line.color.rgb = ACCENT
ci_box.line.width = Pt(1)

ci_eq = render_latex(r"$\mathrm{Conditional\ Independence:}\quad y_i \perp y_{i^\prime} \mid \eta_i$",
                     "eq_ci.png", fontsize=20)
sl.shapes.add_picture(ci_eq, Inches(2.5), Inches(4.65), Inches(7.5), Inches(0.55))

# labels
labels = [
    ("   where:", Inches(5.6)),
    ("   alpha = global intercept    |    beta = K=43 fixed effects    |    u_j = J=19 specialty random intercepts",
     Inches(6.0)),
]
for txt, t in labels:
    add_textbox(sl, Inches(1.0), t, Inches(11), Inches(0.45),
                txt, font_size=17, color=TEXT_DARK)

# =============================================================================
# SLIDE 12: Prior Specification
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Prior Specification")

priors_data = [
    (r"$\alpha \;\sim\; \mathcal{N}(0,\; 5^2)$",
     "Weakly informative on log-odds scale;\ncovers plausible baseline readmission rates"),
    (r"$\beta_k \;\sim\; \mathcal{N}(0,\; 2.5^2)$",
     "Allows OR up to exp(5) ~ 150;\nprevents extreme inflation without over-constraining"),
    (r"$u_j \mid \tau \;\sim\; \mathcal{N}(0,\; \tau^2)$",
     "Exchangeable random intercepts;\nenables partial pooling across specialties"),
    (r"$\tau \;\sim\; \mathrm{Half{-}Normal}(0,\; 1)$",
     "Hyperprior on between-group SD;\ncontrols degree of shrinkage"),
]

top = Inches(1.5)
for i, (tex, desc) in enumerate(priors_data):
    # render formula
    fname = f"eq_prior_{i}.png"
    try:
        img = render_latex(tex, fname, fontsize=22)
        sl.shapes.add_picture(img, Inches(0.8), top, Inches(5.5), Inches(0.6))
    except Exception:
        # fallback: plain text for tau
        labels_fallback = [
            "alpha ~ N(0, 5^2)",
            "beta_k ~ N(0, 2.5^2)",
            "u_j | tau ~ N(0, tau^2)",
            "tau ~ Half-Normal(0, 1)",
        ]
        add_textbox(sl, Inches(0.8), top, Inches(5.5), Inches(0.6),
                    labels_fallback[i], font_size=20, bold=True, color=DARK_BLUE)

    # description
    add_textbox(sl, Inches(6.5), top + Inches(0.02), Inches(6.5), Inches(0.7),
                desc, font_size=15, color=TEXT_GRAY)

    top += Inches(1.2)

# Key insight box
kbox = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.8), top + Inches(0.3), Inches(11.5), Inches(0.8))
kbox.fill.solid()
kbox.fill.fore_color.rgb = LIGHT_BLUE_BG
kbox.line.color.rgb = ACCENT

add_textbox(sl, Inches(1.2), top + Inches(0.4), Inches(10.5), Inches(0.6),
            "The key mechanism: u_j | tau ~ N(0, tau^2) induces partial pooling.\n"
            "Small groups are shrunk toward the population mean; tau controls shrinkage strength.",
            font_size=16, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 13: DAG
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Model Hierarchy (DAG)")

# Create DAG using matplotlib
fig, ax = plt.subplots(figsize=(10, 6))
ax.set_xlim(0, 10)
ax.set_ylim(0, 7)
ax.axis("off")

# Nodes
nodes = {
    "tau":   (5, 6.0, r"$\tau$", "#E3F2FD"),
    "u_j":   (5, 4.2, r"$u_j$", "#E8F5E9"),
    "alpha": (2, 4.2, r"$\alpha$", "#FFF3E0"),
    "beta":  (2, 2.4, r"$\boldsymbol{\beta}$", "#FFF3E0"),
    "x_i":   (3.5, 0.8, r"$\mathbf{x}_i$", "#F3E5F5"),
    "eta_i":  (5, 2.4, r"$\eta_i$", "#FCE4EC"),
    "y_i":   (8, 2.4, r"$y_i$", "#FFEBEE"),
}

for name, (x, y, label, color) in nodes.items():
    circle = plt.Circle((x, y), 0.55, facecolor=color, edgecolor="#1B2A4A",
                         linewidth=2, zorder=3)
    ax.add_patch(circle)
    ax.text(x, y, label, fontsize=20, ha="center", va="center", zorder=4,
            math_fontfamily="cm")

# Arrows
import matplotlib.patches as mpatches
arrow_style = mpatches.FancyArrowPatch
arrows = [
    ("tau", "u_j"),
    ("u_j", "eta_i"),
    ("alpha", "eta_i"),
    ("beta", "eta_i"),
    ("x_i", "eta_i"),
    ("eta_i", "y_i"),
]

for start, end in arrows:
    sx, sy = nodes[start][0], nodes[start][1]
    ex, ey = nodes[end][0], nodes[end][1]
    # Adjust for circle radius
    dx, dy = ex - sx, ey - sy
    dist = np.sqrt(dx**2 + dy**2)
    sx2 = sx + 0.55 * dx / dist
    sy2 = sy + 0.55 * dy / dist
    ex2 = ex - 0.55 * dx / dist
    ey2 = ey - 0.55 * dy / dist
    ax.annotate("", xy=(ex2, ey2), xytext=(sx2, sy2),
                arrowprops=dict(arrowstyle="->", lw=2, color="#1B2A4A"),
                zorder=2)

# Level labels
ax.text(9.2, 6.0, "Hyperparameter\nlevel", fontsize=12, ha="center", va="center",
        color="#666666", style="italic")
ax.text(9.2, 4.2, "Group\nlevel (j=1..19)", fontsize=12, ha="center", va="center",
        color="#666666", style="italic")
ax.text(9.2, 2.4, "Observation\nlevel (i=1..N)", fontsize=12, ha="center", va="center",
        color="#666666", style="italic")

# Plate for j
rect_j = plt.Rectangle((4.0, 3.3), 2.0, 2.0, fill=False,
                         edgecolor="#3A7CA5", linewidth=2, linestyle="--", zorder=1)
ax.add_patch(rect_j)
ax.text(5.8, 3.5, "j = 1..J", fontsize=11, color="#3A7CA5")

# Plate for i
rect_i = plt.Rectangle((1.2, 0.0), 8.0, 3.6, fill=False,
                         edgecolor="#C0392B", linewidth=2, linestyle="--", zorder=1)
ax.add_patch(rect_i)
ax.text(8.8, 0.2, "i = 1..N", fontsize=11, color="#C0392B")

fig.savefig(str(TEMP_DIR / "dag.png"), dpi=200, bbox_inches="tight",
            transparent=True, pad_inches=0.2)
plt.close(fig)

sl.shapes.add_picture(str(TEMP_DIR / "dag.png"),
                       Inches(1.5), Inches(1.3), Inches(10), Inches(5.8))

# =============================================================================
# SLIDE 14: Section — Key Derivation
# =============================================================================
section_divider(prs, "Part 4", "Key Derivation")

# =============================================================================
# SLIDE 15: Joint Posterior
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Joint Posterior Distribution")

eq_post = render_latex(
    r"$p(\alpha,\,\boldsymbol{\beta},\,\mathbf{u},\,\tau \mid \mathbf{y})"
    r"\;\propto\; \prod_{i=1}^{N} p(y_i \mid \eta_i)"
    r"\;\cdot\; p(\alpha)"
    r"\;\cdot\; \prod_{k=1}^{K} p(\beta_k)"
    r"\;\cdot\; \prod_{j=1}^{J} p(u_j \mid \tau)"
    r"\;\cdot\; p(\tau)$",
    "eq_joint_post.png", fontsize=20)
sl.shapes.add_picture(eq_post, Inches(0.5), Inches(1.5), Inches(12), Inches(0.9))

# labels for each term
terms = [
    ("Likelihood", "Logistic link, data-dependent"),
    ("Intercept prior", "N(0, 25)"),
    ("Fixed-effect priors", "N(0, 6.25) each"),
    ("Hierarchical prior", "u_j | tau ~ N(0, tau^2)"),
    ("Hyperprior", "Half-Normal(0, 1)"),
]
top = Inches(2.8)
for label, desc in terms:
    add_textbox(sl, Inches(1.5), top, Inches(3.5), Inches(0.4),
                label, font_size=17, bold=True, color=DARK_BLUE)
    add_textbox(sl, Inches(5.5), top, Inches(7), Inches(0.4),
                desc, font_size=17, color=TEXT_GRAY)
    top += Inches(0.48)

# Non-conjugacy box
nc_box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.2))
nc_box.fill.solid()
nc_box.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
nc_box.line.color.rgb = HIGHLIGHT
nc_box.line.width = Pt(2)

add_textbox(sl, Inches(1.4), Inches(5.55), Inches(10.5), Inches(0.9),
            "The posterior is non-conjugate due to the logistic likelihood.\n"
            "No closed-form solution exists --> This motivates MCMC sampling.",
            font_size=19, bold=True, color=HIGHLIGHT, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 16: Log-Posterior
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Log-Posterior Decomposition")

# Log-posterior equation - split into labeled parts
eq_lp_main = render_latex(
    r"$\ell(\alpha,\boldsymbol{\beta},\mathbf{u},\tau)"
    r"\;=\; \sum_{i=1}^{N}\left[y_i\,\eta_i - \log(1+e^{\eta_i})\right]"
    r"\;-\;\frac{\alpha^2}{2\sigma_\alpha^2}"
    r"\;-\;\sum_{k=1}^{K}\frac{\beta_k^2}{2s_\beta^2}$",
    "eq_logpost1.png", fontsize=20)
sl.shapes.add_picture(eq_lp_main, Inches(0.3), Inches(1.4), Inches(12.5), Inches(0.8))

eq_lp_main2 = render_latex(
    r"$\quad\quad\quad\quad\quad\quad\quad\quad\quad"
    r"\;-\;J\log\tau\;-\;\sum_{j=1}^{J}\frac{u_j^2}{2\tau^2}"
    r"\;-\;\frac{\tau^2}{2s_\tau^2}$",
    "eq_logpost2.png", fontsize=20)
sl.shapes.add_picture(eq_lp_main2, Inches(0.3), Inches(2.3), Inches(12.5), Inches(0.8))

# Term labels in colored boxes
term_info = [
    ("(I)  Log-likelihood", "Data-dependent: logistic log-likelihood", ACCENT),
    ("(II) Prior on alpha", "Gaussian regularization on intercept", TEXT_GRAY),
    ("(III) Prior on beta", "Independent Gaussian on each fixed effect", TEXT_GRAY),
    ("(IV) Prior on u | tau", "Hierarchical: tau controls shrinkage;  tau only here, NOT in (I)", HIGHLIGHT),
    ("(V) Prior on tau", "Half-Normal hyperprior on group-level SD", TEXT_GRAY),
]
top = Inches(3.4)
for label, desc, col in term_info:
    add_textbox(sl, Inches(1.0), top, Inches(3.5), Inches(0.4),
                label, font_size=16, bold=True, color=col)
    add_textbox(sl, Inches(4.8), top, Inches(8), Inches(0.4),
                desc, font_size=15, color=TEXT_DARK)
    top += Inches(0.42)

# Key insight box
insight_box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0))
insight_box.fill.solid()
insight_box.fill.fore_color.rgb = LIGHT_BLUE_BG
insight_box.line.color.rgb = ACCENT
insight_box.line.width = Pt(2)

add_textbox(sl, Inches(1.2), Inches(5.9), Inches(10.7), Inches(0.8),
            "Key insight:  tau appears ONLY in terms (IV) and (V), not in the likelihood.\n"
            "tau is informed by data only indirectly through posterior draws of u_j.\n"
            "This means the tau MCMC update does NOT require likelihood evaluation --> O(J) cost.",
            font_size=15, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 17: Reparameterization
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Reparameterization:  phi = log(tau)")

# Why
add_textbox(sl, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
            "Why reparameterize?", font_size=22, bold=True, color=DARK_BLUE)

why_bullets = [
    "tau > 0  constraint makes random walk on tau inefficient near boundary",
    "phi = log(tau) maps to unconstrained R, enabling standard Gaussian proposals",
    "Requires Jacobian correction:  + phi  (from |d(tau)/d(phi)| = exp(phi))",
]
top = Inches(2.0)
for b in why_bullets:
    add_textbox(sl, Inches(1.2), top, Inches(11), Inches(0.4),
                "   " + b, font_size=17, color=TEXT_DARK)
    top += Inches(0.45)

# Transformed log-target
eq_reparam = render_latex(
    r"$h(\phi) \;=\; -\,J\,\phi \;-\; \frac{1}{2}\sum_{j=1}^{J} u_j^2\,e^{-2\phi}"
    r"\;+\; \phi \;-\; \frac{e^{2\phi}}{2\,s_\tau^2}$",
    "eq_reparam.png", fontsize=22)
sl.shapes.add_picture(eq_reparam, Inches(1.0), Inches(3.8), Inches(10), Inches(0.9))

# label terms
term_labels = [
    ("-J phi :  from  -J log(tau)", Inches(5.0)),
    ("+ phi :  Jacobian correction", Inches(5.45)),
    ("exp(-2 phi) :  tau^{-2}  in random-effect variance", Inches(5.9)),
    ("Does NOT involve the likelihood  -->  O(J) cost per update", Inches(6.35)),
]
for txt, t in term_labels:
    add_textbox(sl, Inches(1.5), t, Inches(10), Inches(0.4),
                "   " + txt, font_size=16, color=TEXT_GRAY)

# =============================================================================
# SLIDE 18: Section — Sampling Strategy
# =============================================================================
section_divider(prs, "Part 5", "Sampling Strategy")

# =============================================================================
# SLIDE 19: MwG Block Structure
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Metropolis-within-Gibbs:  Block Structure")

blocks = [
    ("Block 1:  alpha",
     "Symmetric random walk;  incremental update  delta_eta = alpha* - alpha",
     "O(N)"),
    ("Block 2:  beta_k  (k = 1..43)",
     "Component-wise update;  delta_eta = (beta_k* - beta_k) * X[:,k]",
     "O(N) each"),
    ("Block 3:  u_j  (j = 1..19)",
     "Component-wise;  likelihood uses only group-j observations (n_j << N)",
     "O(n_j) each"),
    ("Block 4:  log(tau)",
     "Does not involve likelihood;  only depends on current u_j draws",
     "O(J)"),
    ("Block 5:  Translation Move",
     "Joint alpha-u update exploiting non-identifiability (see next slide)",
     "O(J) x 10"),
]

top = Inches(1.4)
for label, desc, cost in blocks:
    # box
    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.6), top, Inches(12), Inches(0.95))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1)

    add_textbox(sl, Inches(0.9), top + Inches(0.05), Inches(5), Inches(0.4),
                label, font_size=17, bold=True, color=DARK_BLUE)
    add_textbox(sl, Inches(0.9), top + Inches(0.45), Inches(9), Inches(0.4),
                desc, font_size=14, color=TEXT_GRAY)
    # cost badge
    add_textbox(sl, Inches(10.5), top + Inches(0.25), Inches(2), Inches(0.4),
                cost, font_size=15, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

    top += Inches(1.1)

# Overall complexity
add_textbox(sl, Inches(0.8), top + Inches(0.15), Inches(11), Inches(0.5),
            "Overall:  each iteration scales linearly in N, not N x K",
            font_size=18, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 20: Translation Move
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Joint alpha-u Translation Move")

eq_tm1 = render_latex(
    r"$\mathrm{Propose:}\quad \delta \sim \mathcal{N}(0,\,\sigma^2),"
    r"\quad \alpha^* = \alpha + \delta,\quad u_j^* = u_j - \delta$",
    "eq_tm1.png", fontsize=20)
sl.shapes.add_picture(eq_tm1, Inches(0.8), Inches(1.4), Inches(11), Inches(0.7))

# Key property
eq_tm2 = render_latex(
    r"$\eta_i^* = \alpha^* + \mathbf{x}_i^\top\boldsymbol{\beta} + u_{j[i]}^*"
    r"\;=\; (\alpha+\delta) + \mathbf{x}_i^\top\boldsymbol{\beta} + (u_{j[i]}-\delta)"
    r"\;=\; \eta_i$",
    "eq_tm2.png", fontsize=18)
sl.shapes.add_picture(eq_tm2, Inches(0.8), Inches(2.3), Inches(11), Inches(0.7))

# highlight box
hbox = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.7))
hbox.fill.solid()
hbox.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
hbox.line.color.rgb = GREEN
hbox.line.width = Pt(2)
add_textbox(sl, Inches(1.3), Inches(3.3), Inches(10.7), Inches(0.5),
            "Likelihood invariant!  Acceptance ratio depends only on priors --> O(J) cost",
            font_size=18, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# MH ratio
eq_tm3 = render_latex(
    r"$\log R = \left[-\frac{{\alpha^*}^2}{2\sigma_\alpha^2}"
    r"\;-\;\sum_j\frac{{u_j^*}^2}{2\tau^2}\right]"
    r"\;-\;\left[-\frac{\alpha^2}{2\sigma_\alpha^2}"
    r"\;-\;\sum_j\frac{u_j^2}{2\tau^2}\right]$",
    "eq_tm3.png", fontsize=18)
sl.shapes.add_picture(eq_tm3, Inches(0.8), Inches(4.2), Inches(11), Inches(0.8))

# insights
tm_insights = [
    "Exploits the non-identifiability between alpha and u_j",
    "Repeated 10 times per iteration to maximize mixing benefit",
    "Dramatically improves exploration of the alpha-u ridge",
]
top = Inches(5.3)
for ins in tm_insights:
    add_textbox(sl, Inches(1.2), top, Inches(10.5), Inches(0.4),
                "   " + ins, font_size=17, color=TEXT_DARK)
    top += Inches(0.42)

# =============================================================================
# SLIDE 21: Diagnostics
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "MCMC Settings & Convergence Diagnostics")

settings_text = [
    "4 chains  |  30,000 iterations  |  Burn-in: 8,000  |  Thin by 5  -->  4,400 saved samples / chain",
    "Adaptive per-coordinate proposal tuning:  target acceptance rate 20% - 50%",
]
top = Inches(1.3)
for s in settings_text:
    add_textbox(sl, Inches(0.8), top, Inches(12), Inches(0.4),
                "   " + s, font_size=16, color=TEXT_DARK)
    top += Inches(0.4)

# trace plot figure
fig2 = str(FIG_DIR / "fig2_trace_plots.png")
if os.path.exists(fig2):
    sl.shapes.add_picture(fig2, Inches(0.5), Inches(2.3), Inches(8.5), Inches(4.8))

# diagnostics summary box
diag_box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(9.3), Inches(2.5), Inches(3.7), Inches(4.0))
diag_box.fill.solid()
diag_box.fill.fore_color.rgb = LIGHT_BG
diag_box.line.color.rgb = ACCENT

diag_text = [
    ("Split-Rhat", "All < 1.05"),
    ("ESS (alpha)", "~306"),
    ("ESS (beta, mean)", "~1,200"),
    ("ESS (u_j, mean)", "~3,400"),
    ("ESS (tau)", "~2,412"),
]
dy = Inches(2.6)
add_textbox(sl, Inches(9.5), dy, Inches(3.3), Inches(0.4),
            "Diagnostics Summary", font_size=16, bold=True, color=DARK_BLUE)
dy += Inches(0.5)
for label, val in diag_text:
    add_textbox(sl, Inches(9.5), dy, Inches(2), Inches(0.35),
                label, font_size=14, color=TEXT_GRAY)
    add_textbox(sl, Inches(11.3), dy, Inches(1.5), Inches(0.35),
                val, font_size=14, bold=True, color=GREEN)
    dy += Inches(0.38)

# =============================================================================
# SLIDE 22: Section — Inference
# =============================================================================
section_divider(prs, "Part 6", "Inferential Findings")

# =============================================================================
# SLIDE 23: Fixed Effects
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Fixed Effects — Key Predictors")

add_textbox(sl, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
            "Effects on log-odds scale; interpretable as multiplicative changes in odds (OR = exp(beta))",
            font_size=15, color=TEXT_GRAY)

# Read table2 and pick top predictors (by |post_mean|, excluding alpha)
rows = []
with open(str(TBL_DIR / "table2_fixed_effects.csv")) as f:
    reader = csv.DictReader(f)
    for r in reader:
        if "alpha" in r["parameter"]:
            continue
        rows.append(r)

# Sort by absolute post_mean
rows.sort(key=lambda r: abs(float(r["post_mean"])), reverse=True)
top_rows = rows[:12]

# Table
headers = ["Variable", "Post. Mean", "95% CrI", "OR"]
col_x = [Inches(0.6), Inches(5.8), Inches(7.5), Inches(10.5)]
col_w = [Inches(5.2), Inches(1.7), Inches(3.0), Inches(2.0)]
top = Inches(1.7)

# Header row
for cx, cw, h in zip(col_x, col_w, headers):
    add_textbox(sl, cx, top, cw, Inches(0.35),
                h, font_size=14, bold=True, color=WHITE)

hdr_bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.5), top - Inches(0.02), Inches(12.2), Inches(0.38))
hdr_bg.fill.solid()
hdr_bg.fill.fore_color.rgb = NAVY
hdr_bg.line.fill.background()
sp = hdr_bg._element
sp.getparent().remove(sp)
sl.shapes._spTree.insert(2, sp)

# Re-add headers on top
for cx, cw, h in zip(col_x, col_w, headers):
    add_textbox(sl, cx, top, cw, Inches(0.35),
                h, font_size=14, bold=True, color=WHITE)

top += Inches(0.4)

for i, r in enumerate(top_rows):
    pm = float(r["post_mean"])
    cl = float(r["ci_lower"])
    cu = float(r["ci_upper"])
    or_val = float(r["OR"])

    # clean variable name
    vname = r["parameter"].replace("_z", "").replace("_", " ")

    # highlight if CrI excludes 0
    sig = (cl > 0 and cu > 0) or (cl < 0 and cu < 0)
    row_color = TEXT_DARK if sig else TEXT_GRAY
    bg_color = RGBColor(0xFD, 0xF2, 0xE9) if sig else (LIGHT_BG if i % 2 == 0 else WHITE)

    row_bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), top, Inches(12.2), Inches(0.38))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = bg_color
    row_bg.line.fill.background()
    sp2 = row_bg._element
    sp2.getparent().remove(sp2)
    sl.shapes._spTree.insert(2, sp2)

    vals = [
        vname,
        f"{pm:+.3f}",
        f"[{cl:.3f}, {cu:.3f}]",
        f"{or_val:.3f}",
    ]
    for cx, cw, v in zip(col_x, col_w, vals):
        add_textbox(sl, cx, top, cw, Inches(0.35),
                    v, font_size=13, color=row_color, bold=sig)
    top += Inches(0.38)

add_textbox(sl, Inches(0.6), top + Inches(0.1), Inches(11), Inches(0.3),
            "Highlighted rows: 95% credible interval excludes zero",
            font_size=13, color=HIGHLIGHT)

# =============================================================================
# SLIDE 24: Specialty Random Effects
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Specialty Random Effects & Shrinkage")

fig4 = str(FIG_DIR / "fig4_specialty_random_effects.png")
if os.path.exists(fig4):
    sl.shapes.add_picture(fig4, Inches(0.3), Inches(1.2), Inches(8), Inches(5.8))

# Side text
side_top = Inches(1.5)
side_items = [
    ("Interpretation", "u_j = deviation from\npopulation-average\nlog-odds"),
    ("Highest risk", "Physical Med & Rehab\n(u ~ +0.5)"),
    ("Lowest risk", "Pediatrics, OB/Gyn\n(u ~ -0.5 to -1.0)"),
    ("Shrinkage", "Small-sample groups\npulled toward the\npopulation mean"),
    ("Key insight", "Partial pooling\nstabilizes estimates\nfor rare specialties"),
]

for label, desc in side_items:
    add_textbox(sl, Inches(8.8), side_top, Inches(4), Inches(0.3),
                label, font_size=15, bold=True, color=DARK_BLUE)
    add_textbox(sl, Inches(8.8), side_top + Inches(0.3), Inches(4), Inches(0.8),
                desc, font_size=14, color=TEXT_GRAY)
    side_top += Inches(1.1)

# =============================================================================
# SLIDE 25: Tau — Core Finding
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Heterogeneity Scale tau  —  Core Finding")

# tau posterior figure
fig_tau = str(FIG_DIR / "fig_tau_posterior.png")
if os.path.exists(fig_tau):
    sl.shapes.add_picture(fig_tau, Inches(0.5), Inches(1.3), Inches(7), Inches(4.5))

# Results box
res_box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(8.0), Inches(1.5), Inches(4.8), Inches(3.5))
res_box.fill.solid()
res_box.fill.fore_color.rgb = LIGHT_BLUE_BG
res_box.line.color.rgb = ACCENT

res_items = [
    ("Posterior mean", "0.437"),
    ("Posterior SD", "0.090"),
    ("95% CrI", "[0.298, 0.644]"),
]
ry = Inches(1.7)
add_textbox(sl, Inches(8.3), ry, Inches(4.2), Inches(0.4),
            "tau Posterior Summary", font_size=18, bold=True, color=DARK_BLUE)
ry += Inches(0.5)
for lab, val in res_items:
    add_textbox(sl, Inches(8.3), ry, Inches(2.2), Inches(0.35),
                lab, font_size=16, color=TEXT_GRAY)
    add_textbox(sl, Inches(10.5), ry, Inches(2), Inches(0.35),
                val, font_size=16, bold=True, color=DARK_BLUE)
    ry += Inches(0.4)

# interpretation
interp_box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.0))
interp_box.fill.solid()
interp_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
interp_box.line.color.rgb = GREEN
interp_box.line.width = Pt(2)

add_textbox(sl, Inches(1.2), Inches(6.1), Inches(10.7), Inches(0.8),
            "Posterior mass clearly away from zero --> strong evidence of between-specialty heterogeneity.\n"
            "Answer: Specialty heterogeneity exists and is reliably estimated,\n"
            "but patient-level covariates remain the primary driver of readmission.",
            font_size=16, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 26: PPC & Model Comparison
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Posterior Predictive Check & Model Comparison")

# PPC figure
fig5 = str(FIG_DIR / "fig5_posterior_predictive.png")
if os.path.exists(fig5):
    sl.shapes.add_picture(fig5, Inches(0.3), Inches(1.3), Inches(7.5), Inches(4.5))

# Model comparison table
mc_box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(8.2), Inches(1.3), Inches(4.8), Inches(3.0))
mc_box.fill.solid()
mc_box.fill.fore_color.rgb = LIGHT_BG
mc_box.line.color.rgb = ACCENT

add_textbox(sl, Inches(8.4), Inches(1.4), Inches(4.4), Inches(0.4),
            "Model Comparison", font_size=16, bold=True, color=DARK_BLUE)

mc_data = [
    ("Model", "AUC", "Brier"),
    ("No Specialty", "0.629", "0.085"),
    ("Fixed Effects", "0.641", "0.085"),
    ("Hierarchical", "0.641", "0.085"),
]
my = Inches(1.9)
for i, (m, a, b) in enumerate(mc_data):
    bold_flag = (i == 0)
    col = DARK_BLUE if i == 0 else (GREEN if i == 3 else TEXT_DARK)
    if i == 3:
        bold_flag = True
    add_textbox(sl, Inches(8.4), my, Inches(2.2), Inches(0.3),
                m, font_size=13, bold=bold_flag, color=col)
    add_textbox(sl, Inches(10.5), my, Inches(1), Inches(0.3),
                a, font_size=13, bold=bold_flag, color=col)
    add_textbox(sl, Inches(11.6), my, Inches(1), Inches(0.3),
                b, font_size=13, bold=bold_flag, color=col)
    my += Inches(0.32)

# DIC
add_textbox(sl, Inches(8.4), my + Inches(0.2), Inches(4.4), Inches(0.6),
            "DIC = 23,254\n(Hierarchical < No-specialty)",
            font_size=14, bold=True, color=DARK_BLUE)

# conclusion
conc_box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9))
conc_box.fill.solid()
conc_box.fill.fore_color.rgb = LIGHT_BLUE_BG
conc_box.line.color.rgb = ACCENT

add_textbox(sl, Inches(1.2), Inches(6.1), Inches(10.7), Inches(0.7),
            "Model captures both first-order (mean readmission rate) and second-order (specialty variance) structure.\n"
            "PPC confirms model adequacy: replicated statistics consistent with observed data.",
            font_size=16, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 27: Section — Sensitivity
# =============================================================================
section_divider(prs, "Part 7", "Sensitivity & Robustness Analysis")

# =============================================================================
# SLIDE 28: Sensitivity Analysis
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Sensitivity Analysis — Prior Robustness")

# Prior configs table
add_textbox(sl, Inches(0.6), Inches(1.3), Inches(5), Inches(0.4),
            "Alternative Prior Configurations:", font_size=17, bold=True, color=DARK_BLUE)

configs = [
    ("Baseline", "s_beta = 2.5,  s_tau = 1.0  (Half-Normal)"),
    ("A1", "s_beta = 5.0  (doubled)"),
    ("B1", "s_tau = 2.5  (Half-Normal, wider)"),
    ("B2", "s_tau = 2.5  (Half-Cauchy, heavier tails)"),
]
cy = Inches(1.8)
for label, desc in configs:
    add_textbox(sl, Inches(0.9), cy, Inches(1.2), Inches(0.3),
                label, font_size=14, bold=True, color=ACCENT)
    add_textbox(sl, Inches(2.2), cy, Inches(4.5), Inches(0.3),
                desc, font_size=14, color=TEXT_DARK)
    cy += Inches(0.32)

# sensitivity tau figure
fig6 = str(FIG_DIR / "fig6_sensitivity_tau.png")
if os.path.exists(fig6):
    sl.shapes.add_picture(fig6, Inches(0.3), Inches(3.5), Inches(6.5), Inches(3.5))

# Results box
sens_box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(7.2), Inches(3.5), Inches(5.8), Inches(3.5))
sens_box.fill.solid()
sens_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
sens_box.line.color.rgb = GREEN

sens_results = [
    "tau range across configs:  < 0.05",
    "Beta variation:  < 0.02 for top predictors",
    "Specialty ranking:  Spearman rho > 0.95",
    "",
    "Conclusion:",
    "Posterior is likelihood-dominated",
    "for key parameters.",
    "Results are robust to prior choice.",
]
sy = Inches(3.7)
for i, s in enumerate(sens_results):
    bold_f = i >= 4
    col = GREEN if i >= 4 else TEXT_DARK
    add_textbox(sl, Inches(7.5), sy, Inches(5.2), Inches(0.35),
                s, font_size=15, bold=bold_f, color=col)
    sy += Inches(0.35)

# =============================================================================
# SLIDE 29: Section — Conclusion
# =============================================================================
section_divider(prs, "Part 8", "Conclusion")

# =============================================================================
# SLIDE 30: Summary
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Summary")

summary_items = [
    ("Specialty heterogeneity",
     "tau posterior clearly away from zero (95% CrI [0.30, 0.64])\n"
     "--> Strong evidence of between-specialty variation in readmission risk"),
    ("Patient-level drivers",
     "Prior inpatient visits, comorbid neoplasms, diabetes medication, number of diagnoses\n"
     "are the primary predictors of readmission"),
    ("Partial pooling",
     "Hierarchical model stabilizes estimates for small specialties via shrinkage\n"
     "toward the population mean"),
    ("Model adequacy",
     "PPC confirms model captures both mean and variance structure;\n"
     "DIC favors hierarchical over no-specialty model"),
    ("Robustness",
     "Conclusions stable across alternative prior specifications;\n"
     "posterior is likelihood-dominated"),
]

top = Inches(1.4)
for label, desc in summary_items:
    lbox = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.6), top, Inches(12), Inches(1.0))
    lbox.fill.solid()
    lbox.fill.fore_color.rgb = LIGHT_BG
    lbox.line.color.rgb = ACCENT
    lbox.line.width = Pt(1)

    add_textbox(sl, Inches(0.9), top + Inches(0.05), Inches(4), Inches(0.35),
                label, font_size=17, bold=True, color=DARK_BLUE)
    add_textbox(sl, Inches(0.9), top + Inches(0.4), Inches(11.5), Inches(0.55),
                desc, font_size=14, color=TEXT_DARK)
    top += Inches(1.15)

# =============================================================================
# SLIDE 31: Limitations & Future Work
# =============================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(sl, "Limitations & Future Work")

lim_items = [
    ("Data vintage", "Dataset from 1999-2008;  ICD-9 has since transitioned to ICD-10;\n"
     "clinical practices may have evolved"),
    ("Selection bias", "Excluded 45.7% of encounters with unknown specialty;\n"
     "remaining cohort may not be fully representative"),
    ("Model scope", "Random intercepts only (no random slopes);\n"
     "assumes specialty effect is constant across patient profiles"),
]

fut_items = [
    "Add temporal effects (year/season) to capture time trends",
    "Incorporate hospital-level hierarchy (patients nested in hospitals nested in specialties)",
    "Extend to random slopes for key covariates (e.g., age x specialty interaction)",
    "Apply to modern datasets with ICD-10 coding and broader demographic coverage",
]

top = Inches(1.5)
add_textbox(sl, Inches(0.8), top, Inches(5), Inches(0.4),
            "Limitations", font_size=20, bold=True, color=HIGHLIGHT)
top += Inches(0.5)
for label, desc in lim_items:
    add_textbox(sl, Inches(1.0), top, Inches(2.5), Inches(0.35),
                label, font_size=16, bold=True, color=TEXT_DARK)
    add_textbox(sl, Inches(3.5), top, Inches(9), Inches(0.6),
                desc, font_size=15, color=TEXT_GRAY)
    top += Inches(0.7)

top += Inches(0.3)
add_textbox(sl, Inches(0.8), top, Inches(5), Inches(0.4),
            "Future Directions", font_size=20, bold=True, color=GREEN)
top += Inches(0.5)
for f in fut_items:
    add_textbox(sl, Inches(1.0), top, Inches(11.5), Inches(0.4),
                "   " + f, font_size=16, color=TEXT_DARK)
    top += Inches(0.45)

# Thank you
top += Inches(0.3)
add_textbox(sl, Inches(0.8), top, Inches(12), Inches(0.5),
            "Thank you!  Questions?",
            font_size=28, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

# ── save ──────────────────────────────────────────────────────────────────────
prs.save(str(OUT_PPT))
print(f"Presentation saved to: {OUT_PPT}")
print(f"Total slides: {len(prs.slides)}")
